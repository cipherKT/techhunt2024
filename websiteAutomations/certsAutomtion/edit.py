from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches
import json

OUTPUT_DIR = "./certs"
standings = ["FIRST", "SECOND", "THIRD", "FOURTH", "FIFTH", "SIXTH", "SEVENTH", "EIGHTH", "NINTH", "TENTH"]

def filename(team_id, name):
    return f"{team_id}_{name}"

def putName(shape, user):
    text_frame = shape.text_frame
    text_frame.clear()
    paragraph = text_frame.add_paragraph()
    name = user["name"].upper()
    paragraph.text = name[:17] + '...' if len(name) > 17 else name
    
    shape.top = Inches(1.80)
    
    run = paragraph.runs[0]
    run.font.name = 'Inter Black'
    run.font.bold = True
    run.font.size = 300000
    run.font.color.rgb = RGBColor(59, 130, 246)
    
def putTeam(shape, user):
    text_frame = shape.text_frame
    text_frame.clear()
    paragraph = text_frame.add_paragraph()
    paragraph.text = user["team_name"].upper()
    
    shape.top = Inches(2.51)
    
    run = paragraph.runs[0]
    run.font.name = 'Anek Bangla SemiBold'
    run.font.bold = True
    run.font.size = 150000
    run.font.color.rgb = RGBColor(185, 185, 185)

def putPosition(shape, user):
    text_frame = shape.text_frame
    text_frame.clear()
    paragraph = text_frame.add_paragraph()
    position = standings[user["standing"] - 1].upper() + " POSITION"
    paragraph.text = position
    
    shape.top = Inches(2.65)
    
    run = paragraph.runs[0]
    run.font.bold = True
    run.font.size = 180000
    run.font.name = 'Inter'
    run.font.color.rgb = RGBColor(190, 190, 190)

if __name__ == "__main__":
    with open("standings.json") as file:
        users = json.load(file)
    
    for user in users:
        if user["standing"] <= 10: 
            presentation = Presentation("Cert_Top.pptx")
        else:
            presentation = Presentation("Cert_Bottom.pptx")
            
        slide = presentation.slides[0]
        
        for shape in slide.shapes: 
            if shape.has_text_frame:
                if shape.text == "[NAME]":
                    putName(shape, user)
                elif shape.text == "[TEAM NAME]":
                    putTeam(shape, user)
                if user["standing"] <= 10 and shape.text == "[POSITION]":
                    putPosition(shape, user)
        
        pic = slide.shapes.add_picture(
            f"./qrs/{filename(user["team_id"], user["name"])}.png", 
            Inches(7), 
            Inches(4.45), 
            Inches(0.9015748), 
            Inches(0.9015748)
        )
         
        if user["standing"] <= 10: 
            presentation.save(OUTPUT_DIR + f"/winners/{filename(user["team_id"], user["name"])}.pptx")
        else:
            presentation.save(OUTPUT_DIR + f"/loosers/{filename(user["team_id"], user["name"])}.pptx")